from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "University_Event_Management_System_Presentation.pptx"

TITLE = "University Event Management System"
SUBTITLE = "Project Presentation"

BG = RGBColor(248, 246, 242)
NAVY = RGBColor(20, 33, 61)
RED = RGBColor(218, 68, 83)
TEXT = RGBColor(35, 35, 35)
MUTED = RGBColor(95, 95, 95)
LIGHT = RGBColor(255, 255, 255)
PLACEHOLDER = RGBColor(220, 220, 220)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.03), Inches(1.3), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.45))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.name = "Aptos"
        sub_p.font.size = Pt(11)
        sub_p.font.color.rgb = MUTED


def add_footer(slide, page_num):
    footer = slide.shapes.add_textbox(Inches(12.6), Inches(7.0), Inches(0.4), Inches(0.3))
    p = footer.text_frame.paragraphs[0]
    p.text = str(page_num)
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5.5), Inches(5.1))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_two_image_placeholders(slide):
    positions = [
        (Inches(7.0), Inches(1.7), Inches(5.6), Inches(2.2), "Screenshot Space 1"),
        (Inches(7.0), Inches(4.25), Inches(5.6), Inches(2.2), "Screenshot Space 2"),
    ]
    for left, top, width, height, label in positions:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = PLACEHOLDER
        shape.line.dash_style = 1
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = MUTED
        p2 = tf.add_paragraph()
        p2.text = "Insert app image here"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = "Aptos"
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.5), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = "Aptos Display"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.0), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = SUBTITLE
    sp.font.name = "Aptos"
    sp.font.size = Pt(18)
    sp.font.color.rgb = RED

    info = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(4.0), Inches(2.0))
    tf = info.text_frame
    for idx, line in enumerate(
        [
            "Prepared for project presentation",
            "Frontend: React + Vite",
            "Backend: Express + TypeScript + Prisma",
            "Database: PostgreSQL",
        ]
    ):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(18 if idx == 0 else 16)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)

    hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.1), Inches(5.4), Inches(5.2))
    hero.fill.solid()
    hero.fill.fore_color.rgb = LIGHT
    hero.line.color.rgb = PLACEHOLDER
    hero.line.dash_style = 1
    tf = hero.text_frame
    p = tf.paragraphs[0]
    p.text = "Cover Screenshot Space"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = MUTED

    add_footer(slide, 1)


def add_content_slide(prs, page_num, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, subtitle)
    add_bullets(slide, bullets)
    add_two_image_placeholders(slide)
    add_footer(slide, page_num)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    slides = [
        ("Agenda", "Topics covered in the presentation", [
            "Introduction to the project and core objective",
            "Problem statement, scope, and target users",
            "System architecture and technology stack",
            "Database design, modules, and workflow explanation",
            "Frontend, backend, deployment, testing, and future scope",
        ]),
        ("Project Overview", "What the system is designed to achieve", [
            "The University Event Management System is a centralized platform for managing university events, registrations, and user roles.",
            "It provides a modern public website for event discovery and a secure admin workflow for event control.",
            "The platform supports event groups, individual events, user accounts, and controlled participation management.",
            "The project aims to reduce manual effort and provide a structured digital process for event operations.",
        ]),
        ("Problem Statement", "Why this system is needed", [
            "University event information is often scattered across posters, social media, and informal communication channels.",
            "Manual registration handling creates delays, duplicate work, and difficulty in tracking participants.",
            "Organizers need a better way to control event visibility, user eligibility, and registration status.",
            "The system solves these issues through one consistent platform with role-based workflows.",
        ]),
        ("Objectives of the Project", "Main goals of development", [
            "Create a single platform to publish and manage university events.",
            "Provide secure student registration with email verification and login support.",
            "Allow admins to manage events, users, and university badge approvals.",
            "Support both solo and team-based participation with proper validation rules.",
            "Prepare a scalable base for future competition and result-management features.",
        ]),
        ("Users and Roles", "Who interacts with the system", [
            "Student: registers, logs in, updates profile, verifies university details, and joins events.",
            "Admin: manages users, event groups, events, approvals, and coordinator assignments.",
            "Coordinator: receives event-specific temporary permissions for operational support.",
            "Public visitor: browses published event content without needing admin access.",
        ]),
        ("Technology Stack", "Tools and frameworks used", [
            "Frontend built with React and Vite for a fast and modern UI experience.",
            "Backend developed using Express and TypeScript for structured API development.",
            "Database handled through PostgreSQL with Prisma ORM for schema-driven access.",
            "Authentication uses JWT, password hashing, and OTP verification workflows.",
            "Cloudinary is used for media upload handling and image hosting.",
        ]),
        ("System Architecture", "High-level structure of the application", [
            "The application follows a layered architecture: frontend, backend API, and database.",
            "The frontend manages screens, forms, navigation, and user interactions.",
            "The backend applies authentication, validation, and business rules before storing data.",
            "The database persists users, events, teams, registrations, and related operational records.",
        ]),
        ("Frontend Module", "User interface and interaction layer", [
            "The frontend includes pages for home, events, schedule, contact, login, signup, profile, and admin dashboard.",
            "React Router is used for navigation across public and protected pages.",
            "Axios handles API communication, and the token is stored in local storage for authenticated requests.",
            "The interface is designed to be responsive and visually suitable for an event-driven application.",
        ]),
        ("Backend Module", "API and business logic layer", [
            "The backend is organized into routes, controllers, services, middleware, validators, and utilities.",
            "Each request follows a clean flow: route to controller to service to Prisma.",
            "This separation improves maintainability and makes business logic easier to update later.",
            "The backend also exposes dedicated routes for auth, profile, admin, events, uploads, and coordinator features.",
        ]),
        ("Database Design", "Core entities used in the project", [
            "The database includes tables and models for users, OTP codes, event groups, events, registrations, teams, and badge logs.",
            "Event groups act as parent collections for related events.",
            "Registrations store event participation state, payment state, and review details.",
            "Team and team member models support multi-user participation in group events.",
        ]),
        ("Authentication and Security", "How access is controlled", [
            "Students register using email and password, then verify their account using OTP.",
            "JWT tokens are used to authorize protected API requests.",
            "Middleware checks user identity and role before allowing admin or protected actions.",
            "University-only access is controlled through badge verification status.",
            "Password reset flow is also secured with OTP validation.",
        ]),
        ("Event Management", "How admins manage content", [
            "Admins can create event groups to represent festivals or larger event collections.",
            "Within each group, child events can be created with details like venue, dates, fees, participation type, and audience scope.",
            "Visibility rules ensure draft or restricted events are handled correctly.",
            "This structure makes the system practical for real campus festival management.",
        ]),
        ("Registration Workflow", "How participants join events", [
            "Students can view published events and register based on event rules.",
            "The backend checks registration windows, audience scope, and capacity limits.",
            "Events can require manual approval or allow direct confirmation.",
            "Payment fields are supported so the platform can later integrate real payment flows.",
        ]),
        ("Team Participation", "Support for team-based events", [
            "The platform supports both solo and team event types.",
            "A team can be created with a captain and multiple members within allowed size limits.",
            "Validation prevents duplicate or invalid team participation.",
            "This module is important for technical and cultural competitions involving groups.",
        ]),
        ("Admin Dashboard", "Operational control panel", [
            "The admin dashboard provides a dedicated workspace for managing platform data.",
            "Admins can create events, assign coordinators, verify badges, and create other admin users.",
            "Counts and metrics help admins track published content and pending actions.",
            "The dashboard is designed for practical use rather than only academic demonstration.",
        ]),
        ("Profile and Badge Verification", "University identity validation", [
            "Users can update academic information such as course, department, year, and university details.",
            "Admins review university badge submissions and mark them as verified or rejected.",
            "Verified badge status enables access to university-only events.",
            "This adds trust, control, and event eligibility filtering to the platform.",
        ]),
        ("Deployment Architecture", "How the project is hosted", [
            "The frontend is deployed separately from the backend for cleaner scaling and simpler hosting.",
            "React + Vite frontend can be deployed on Vercel as a static client application.",
            "Express backend is deployed on Render and connects to PostgreSQL.",
            "Environment variables are used to connect frontend and backend securely across platforms.",
        ]),
        ("Testing and Validation", "How the project is verified", [
            "TypeScript build checks help catch structural errors in the frontend and backend code.",
            "API validation is enforced using request validators and business rule checks.",
            "Manual testing is important for auth flows, event creation, registration, and admin actions.",
            "Deployment testing ensures that Vercel frontend and Render backend communicate correctly.",
        ]),
        ("Future Scope", "Possible enhancements", [
            "Add payment gateway integration for paid events.",
            "Introduce analytics dashboards and downloadable reports.",
            "Expand competition modules for rounds, match progression, and final result publishing.",
            "Add email notifications, attendance tracking, and certificate generation.",
            "Create a mobile-friendly or app-based extension in the future.",
        ]),
        ("Conclusion", "Final summary of the project", [
            "The University Event Management System provides a structured digital solution for event publication and participation management.",
            "It combines a modern frontend, secure backend, and relational database design into one complete project.",
            "The project is practical, scalable, and suitable as both an academic submission and a real-world foundation.",
            "With screenshots and live deployment references added, this presentation will fully explain the project flow.",
        ]),
    ]

    for idx, (title, subtitle, bullets) in enumerate(slides, start=2):
        add_content_slide(prs, idx, title, subtitle, bullets)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
